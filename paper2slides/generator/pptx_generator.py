"""
PPTX Generator — converts generated slide/poster images into an editable .pptx file.
Each image occupies a full-bleed slide with no additional chrome.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)


def create_pptx_from_images(
    image_paths: List[Path],
    output_path: Path,
) -> Path:
    """
    Build a Presentation from *image_paths* (one slide per image).

    The slide dimensions are inferred from the first image so the image
    always fills the slide without letterboxing.

    Args:
        image_paths: Ordered list of image files (png / jpg / webp).
        output_path: Destination .pptx file path.

    Returns:
        The resolved *output_path*.
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Emu  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-pptx is required for PPTX export. Run: pip install python-pptx"
        ) from exc

    try:
        from PIL import Image as _PILImage  # type: ignore[import]

        _pil_available = True
    except ImportError:
        _pil_available = False

    if not image_paths:
        raise ValueError("image_paths must not be empty")

    # ------------------------------------------------------------------
    # Determine slide dimensions from the first image
    # ------------------------------------------------------------------
    first = image_paths[0]
    if _pil_available:
        with _PILImage.open(first) as im:
            img_w_px, img_h_px = im.size
    else:
        # Fallback: assume square (common for Paper2Slides outputs)
        img_w_px, img_h_px = 1024, 1024

    # Scale so width = 10 inches (914400 EMU/inch)
    EMU_PER_INCH = 914_400
    slide_w_emu = 10 * EMU_PER_INCH
    slide_h_emu = int(slide_w_emu * img_h_px / img_w_px)

    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)

    blank_layout = prs.slide_layouts[6]  # index 6 = completely blank

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("Saved PPTX → %s  (%d slide(s))", output_path.name, len(image_paths))
    return output_path
